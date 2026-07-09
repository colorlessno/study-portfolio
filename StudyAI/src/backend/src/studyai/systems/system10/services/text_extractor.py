from __future__ import annotations

from io import BytesIO
from pathlib import Path
from xml.etree import ElementTree
from zipfile import ZipFile

from studyai.common.errors.models import ValidationAppError


class TextExtractor:
    TEXT_EXTENSIONS = {
        ".txt",
        ".md",
        ".py",
        ".js",
        ".ts",
        ".java",
        ".sql",
        ".sh",
        ".yaml",
        ".yml",
        ".json",
        ".xml",
    }

    def extract_text(self, file_name: str, file_bytes: bytes) -> str:
        suffix = Path(file_name).suffix.lower()
        if suffix == ".pdf":
            return self._extract_pdf(file_bytes)
        if suffix == ".docx":
            return self._extract_docx(file_bytes)
        if suffix == ".xlsx":
            return self._extract_xlsx(file_bytes)
        if suffix == ".pptx":
            return self._extract_pptx(file_bytes)
        if suffix in self.TEXT_EXTENSIONS:
            return self._decode_text(file_bytes)
        raise ValidationAppError("unsupported_file_type", "対応していないファイル形式です。")

    def _extract_pdf(self, file_bytes: bytes) -> str:
        try:
            import fitz
        except ImportError as exc:
            raise ValidationAppError("pdf_support_missing", "PDF の読み取りライブラリが利用できません。") from exc
        document = fitz.open(stream=file_bytes, filetype="pdf")
        try:
            texts = [page.get_text("text") for page in document]
        finally:
            document.close()
        return self._normalize_text("\n".join(texts))

    def _extract_docx(self, file_bytes: bytes) -> str:
        with ZipFile(BytesIO(file_bytes)) as archive:
            try:
                xml_bytes = archive.read("word/document.xml")
            except KeyError as exc:
                raise ValidationAppError("invalid_document_file", "DOCX 本文を読み取れません。") from exc
        root = ElementTree.fromstring(xml_bytes)
        namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        texts = [node.text for node in root.findall(".//w:t", namespace) if node.text]
        return self._normalize_text("\n".join(texts))

    def _extract_xlsx(self, file_bytes: bytes) -> str:
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise ValidationAppError("xlsx_support_missing", "Excel の読み取りライブラリが利用できません。") from exc
        workbook = load_workbook(BytesIO(file_bytes), read_only=True, data_only=True)
        lines: list[str] = []
        try:
            for sheet in workbook.worksheets:
                lines.append(f"[sheet] {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                    if values:
                        lines.append(" | ".join(values))
        finally:
            workbook.close()
        return self._normalize_text("\n".join(lines))

    def _extract_pptx(self, file_bytes: bytes) -> str:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ValidationAppError("pptx_support_missing", "PowerPoint の読み取りライブラリが利用できません。") from exc
        presentation = Presentation(BytesIO(file_bytes))
        lines: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            lines.append(f"[slide] {slide_index}")
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and str(text).strip():
                    lines.append(str(text).strip())
        return self._normalize_text("\n".join(lines))

    def _decode_text(self, file_bytes: bytes) -> str:
        for encoding in ("utf-8-sig", "utf-8", "cp932"):
            try:
                return self._normalize_text(file_bytes.decode(encoding))
            except UnicodeDecodeError:
                continue
        raise ValidationAppError("invalid_document_encoding", "テキストファイルの文字コードを判定できません。")

    @staticmethod
    def _normalize_text(text: str) -> str:
        return "\n".join(
            line.rstrip() for line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n")
        ).strip()
